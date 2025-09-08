"""
Slide extraction and conversion module for SlideSpeaker.

This module handles the extraction of content from presentation files (PDF, PPTX, PPT)
and conversion of slides to image formats. It provides robust error handling and
fallback mechanisms for various file formats and conversion scenarios.
"""

import contextlib
import io
import os
import subprocess
from pathlib import Path
from typing import Any

import PyPDF2
from PIL import Image
from pptx import Presentation


class SlideExtractor:
    """Extractor for presentation slides and converter to images"""

    async def extract_slides(self, file_path: Path, file_ext: str) -> list[str]:
        """
        Extract text content from presentation slides.

        This method supports PDF, PPTX, and PPT file formats. It extracts
        text content from each slide for use in script generation.
        """
        if file_ext == ".pdf":
            return await self._extract_pdf_slides(file_path)
        elif file_ext in [".pptx", ".ppt"]:
            return await self._extract_pptx_slides(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

    async def _extract_pdf_slides(self, file_path: Path) -> list[str]:
        """Extract text content from PDF slides"""
        slides = []
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                slides.append(text.strip())
        return slides

    async def _extract_pptx_slides(self, file_path: Path) -> list[str]:
        """Extract text content from PowerPoint slides"""
        slides = []
        try:
            presentation = Presentation(str(file_path))
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                slides.append("\n".join(slide_text))
            # Successfully extracted slides - no need to print
            pass
        except Exception:
            # Error will be logged by caller
            raise
        return slides

    async def convert_to_image(
        self, file_path: Path, file_ext: str, slide_index: int, output_path: Path
    ) -> None:
        """
        Convert a specific slide to an image file.

        This method handles conversion of both PDF and PowerPoint slides to PNG images.
        It uses external tools (pdftoppm, LibreOffice) for conversion with fallback
        mechanisms for error handling.
        """
        if file_ext == ".pdf":
            await self._convert_pdf_to_image(file_path, slide_index, output_path)
        elif file_ext in [".pptx", ".ppt"]:
            await self._convert_pptx_to_image(file_path, slide_index, output_path)

    async def _convert_pdf_to_image(
        self, file_path: Path, page_index: int, output_path: Path
    ) -> None:
        """
        Convert a specific PDF page to a PNG image.

        Uses the pdftoppm utility for high-quality PDF to image conversion.
        Includes timeout handling and fallback to placeholder images on failure.
        """
        # Try to convert PDF page to image using pdftoppm directly
        try:
            # Use pdftoppm to convert PDF page to PNG
            cmd = [
                "pdftoppm",
                "-png",  # Output format
                "-f",
                str(page_index + 1),  # First page
                "-l",
                str(page_index + 1),  # Last page
                "-r",
                "150",  # Resolution DPI
                "-singlefile",  # Only output the specified page
                str(file_path),
                str(output_path.with_suffix("")),  # Remove .png suffix for pdftoppm
            ]

            # Run the command
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)

            if result.returncode == 0:
                # pdftoppm creates the file with the page number appended, we need to rename it
                # Since we use -singlefile, it should create file without page number
                generated_file = str(output_path.with_suffix("")) + ".png"
                if os.path.exists(generated_file):
                    os.rename(generated_file, str(output_path))
                    # Successfully converted - no need to print
                    pass
                else:
                    # Fallback to content-based image if file wasn't created
                    await self._create_pdf_content_image(
                        file_path, page_index, output_path
                    )
            else:
                # Error details logged by caller
                # Fallback to content-based image if conversion fails
                await self._create_pdf_content_image(file_path, page_index, output_path)
        except subprocess.TimeoutExpired:
            # Timeout handled gracefully
            # Fallback to content-based image if conversion times out
            await self._create_pdf_content_image(file_path, page_index, output_path)
        except Exception:
            # Error details logged by caller
            # Fallback to content-based image if conversion fails
            await self._create_pdf_content_image(file_path, page_index, output_path)

    async def _create_pdf_content_image(
        self, file_path: Path, page_index: int, output_path: Path
    ) -> None:
        """
        Create a visually appealing slide image based on PDF page content.

        Generates a well-designed slide representation of the PDF page's text content
        when direct image extraction is not possible or desired.
        """
        try:
            import textwrap

            from PIL import Image, ImageDraw, ImageFont

            # Extract text content from the specific PDF page
            page_text = ""
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                if page_index < len(pdf_reader.pages):
                    page = pdf_reader.pages[page_index]
                    page_text = page.extract_text().strip()

            # Create a 16:9 slide image (1920x1080) - standard video resolution
            width, height = 1920, 1080
            img = Image.new("RGB", (width, height), "#ffffff")
            draw = ImageDraw.Draw(img)

            # Try to load system fonts, fallback to default if not available
            title_font: ImageFont.FreeTypeFont | ImageFont.ImageFont
            content_font: ImageFont.FreeTypeFont | ImageFont.ImageFont
            small_font: ImageFont.FreeTypeFont | ImageFont.ImageFont
            try:
                # Try common system fonts
                title_font = ImageFont.truetype("Arial.ttf", 48)
                content_font = ImageFont.truetype("Arial.ttf", 32)
                small_font = ImageFont.truetype("Arial.ttf", 24)
            except Exception:
                try:
                    # Try alternative font names
                    title_font = ImageFont.truetype("DejaVuSans.ttf", 48)
                    content_font = ImageFont.truetype("DejaVuSans.ttf", 32)
                    small_font = ImageFont.truetype("DejaVuSans.ttf", 24)
                except Exception:
                    # Fallback to default font
                    title_font = ImageFont.load_default()
                    content_font = ImageFont.load_default()
                    small_font = ImageFont.load_default()

            # Add a subtle gradient background
            gradient = Image.new("RGBA", (width, height), (0, 0, 0, 0))
            gradient_draw = ImageDraw.Draw(gradient)

            # Create a vertical gradient from light blue to white
            for y in range(height):
                # Calculate color intensity (0 at top, 200 at bottom)
                intensity = int(200 * (y / height))
                color = (240, 245, 255, 255 - intensity)  # Light blue to transparent
                gradient_draw.line([(0, y), (width, y)], fill=color)

            img = Image.alpha_composite(img.convert("RGBA"), gradient)
            draw = ImageDraw.Draw(img)

            # Add decorative elements
            # Draw a subtle horizontal line
            draw.line([(100, 150), (width - 100, 150)], fill="#4287f5", width=2)

            # Add slide title (Page X)
            title = f"Page {page_index + 1}"
            bbox = draw.textbbox((0, 0), title, font=title_font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            draw.text((x, 180), title, fill="#2c3e50", font=title_font)

            # Add page content
            if page_text:
                # Clean and limit the text content
                lines = page_text.split("\n")
                # Remove empty lines and limit to first 20 lines
                content_lines = [line.strip() for line in lines if line.strip()][:20]

                if content_lines:
                    # Start content below title
                    content_y = 280

                    # Process content in chunks to fit on slide
                    current_y = content_y
                    for i, line in enumerate(content_lines):
                        if current_y > height - 100:  # Don't go beyond image bounds
                            break

                        # Skip very short lines that are likely just formatting
                        if len(line) < 3:
                            continue

                        # Wrap long lines
                        wrapped_lines = textwrap.wrap(line, width=70)

                        for wrapped_line in wrapped_lines:
                            if current_y > height - 100:  # Don't go beyond image bounds
                                break

                            # Use different font sizes for headers vs content
                            if (
                                i == 0
                                and len(wrapped_line) < 50
                                and wrapped_line.endswith(":")
                            ):
                                # Likely a header
                                font = title_font
                                line_height = 60
                            elif len(wrapped_line) < 30:
                                # Short line, use smaller font
                                font = small_font
                                line_height = 30
                            else:
                                # Regular content
                                font = content_font
                                line_height = 40

                            draw.text(
                                (100, current_y),
                                wrapped_line,
                                fill="#34495e",
                                font=font,
                            )
                            current_y += line_height
            else:
                # No text content, show a message
                message = "No text content available for this page"
                bbox = draw.textbbox((0, 0), message, font=content_font)
                text_width = bbox[2] - bbox[0]
                x = (width - text_width) // 2
                draw.text((x, height // 2), message, fill="#7f8c8d", font=content_font)

            # Add a subtle decorative element in the corner
            # Draw a small circle in bottom right corner
            circle_x, circle_y = width - 50, height - 50
            draw.ellipse(
                [circle_x - 15, circle_y - 15, circle_x + 15, circle_y + 15],
                fill="#4287f5",
                outline="#2c3e50",
                width=1,
            )

            # Add slide number in bottom center
            slide_number_text = f"{page_index + 1}"
            bbox = draw.textbbox((0, 0), slide_number_text, font=small_font)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            draw.text(
                (x, height - 40), slide_number_text, fill="#7f8c8d", font=small_font
            )

            # Save the image
            img.convert("RGB").save(output_path)

        except Exception:
            # Final fallback to simple placeholder
            self._create_placeholder_image(page_index, output_path)

    def _create_placeholder_image(self, slide_index: int, output_path: Path) -> None:
        """Create a placeholder image with slide number"""
        from PIL import Image, ImageDraw

        # Create a 16:9 slide image (1920x1080) - standard video resolution
        width, height = 1920, 1080
        img = Image.new("RGB", (width, height), color="white")
        d = ImageDraw.Draw(img)
        d.text(
            (width // 2, height // 2),
            f"Slide {slide_index + 1}",
            fill="black",
            anchor="mm",
        )
        img.save(output_path)

    async def _convert_pptx_to_image(
        self, file_path: Path, slide_index: int, output_path: Path
    ) -> None:
        """
        Convert a specific PowerPoint slide to a PNG image.

        First converts PPTX to PDF using LibreOffice, then converts the PDF page to image.
        Includes timeout handling and fallback mechanisms for conversion failures.
        """
        # Convert PPTX to PDF first using LibreOffice (soffice), then export PDF pages to images
        try:
            # Convert PPTX to PDF using LibreOffice
            pdf_path = file_path.with_suffix(".pdf")

            # Use LibreOffice to convert PPTX to PDF
            cmd = [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(file_path.parent),
                str(file_path),
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

            if result.returncode == 0 and pdf_path.exists():
                # Successfully converted to PDF, now convert the specific page to image
                await self._convert_pdf_to_image(pdf_path, slide_index, output_path)

                # Clean up the temporary PDF file
                with contextlib.suppress(Exception):
                    pdf_path.unlink()
            else:
                # Conversion failure handled gracefully
                # Fallback to original content-based approach
                await self._convert_pptx_to_image_original(
                    file_path, slide_index, output_path
                )

        except subprocess.TimeoutExpired:
            # Timeout handled gracefully
            # Fallback to original content-based approach
            await self._convert_pptx_to_image_original(
                file_path, slide_index, output_path
            )
        except Exception:
            # Error details logged by caller
            # Fallback to original content-based approach
            await self._convert_pptx_to_image_original(
                file_path, slide_index, output_path
            )

    async def _convert_pptx_to_image_original(
        self, file_path: Path, slide_index: int, output_path: Path
    ) -> None:
        """
        Fallback method for converting PowerPoint slides to images.

        Extracts content directly from the PPTX file and creates an image representation
        based on the slide's text and image content.
        """
        # Original content-based approach as fallback
        try:
            presentation = Presentation(str(file_path))

            # Check if the slide index is valid
            if slide_index >= len(presentation.slides):
                self._create_placeholder_image(slide_index, output_path)
                return

            slide = presentation.slides[slide_index]

            # Try to find images in the slide
            images_found = []
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image:
                    images_found.append(shape.image)

            # If we found images, use the first one
            if images_found:
                image = images_found[0]  # Use the first image found
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                # Resize to standard size if needed
                img = img.resize((800, 600))  # type: ignore
                img.save(output_path, "PNG")
                img.close()  # Close the image to free resources
            else:
                # No images found, create content-based image instead of simple placeholder
                await self._create_content_image(slide, slide_index, output_path)

        except Exception:
            # Error details logged by caller
            # Fallback to placeholder if conversion fails
            self._create_placeholder_image(slide_index, output_path)

    async def _create_content_image(
        self, slide: Any, slide_index: int, output_path: Path
    ) -> None:
        """
        Create an image based on slide content (text, etc.).

        Generates a visually appealing image representation of the slide's text content
        when direct image extraction is not possible.
        """
        try:
            import textwrap

            from PIL import Image, ImageDraw, ImageFont

            # Create a 16:9 slide image (1920x1080) - standard video resolution
            width, height = 1920, 1080
            img = Image.new("RGB", (width, height), color="white")
            draw = ImageDraw.Draw(img)

            # Try to get a font (fallback to default if needed)
            try:
                font = ImageFont.truetype("Arial.ttf", 24)
                title_font = ImageFont.truetype("Arial.ttf", 32)
            except Exception:
                font = ImageFont.load_default()  # type: ignore
                title_font = ImageFont.load_default()  # type: ignore

            # Extract text content from slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                # Draw title (first text element, typically)
                if slide_text:
                    title = slide_text[0]
                    # Wrap title text
                    wrapped_title = textwrap.fill(title, width=40)
                    draw.text((50, 50), wrapped_title, fill="black", font=title_font)

                # Draw other content
                y_position = 120
                for text in slide_text[1:]:
                    if (
                        text and y_position < height - 50
                    ):  # Don't go beyond image bounds
                        wrapped_text = textwrap.fill(text, width=60)
                        draw.text(
                            (50, y_position), wrapped_text, fill="black", font=font
                        )
                        # Calculate text height for next position
                        bbox = draw.textbbox((50, y_position), wrapped_text, font=font)
                        text_height = bbox[3] - bbox[1]
                        y_position += int(text_height) + 10  # Add some spacing
            else:
                # No text content, create simple placeholder
                draw.text(
                    (width // 2, height // 2),
                    f"Slide {slide_index + 1}",
                    fill="black",
                    anchor="mm",
                    font=font,
                )

            img.save(output_path, "PNG")
            # Successfully created content image

        except Exception:
            # Final fallback to simple placeholder
            self._create_placeholder_image(slide_index, output_path)
